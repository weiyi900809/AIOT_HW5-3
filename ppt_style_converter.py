# PPT 風格自動重新設計工具 (PPT Style AI Redesign Tool)
# 支援將任意 PPT 轉換為多種不同的設計風格
# 使用 python-pptx 庫進行程式化修改

"""
使用方法:
    python ppt_style_converter.py input.pptx --styles modern minimal corporate

安裝依賴:
    pip install python-pptx pillow
"""

import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import argparse
from typing import List, Dict, Tuple
from dataclasses import dataclass
from datetime import datetime


@dataclass
class StylePreset:
    """設計風格預設配置"""
    name: str
    description: str
    # 色彩方案
    primary_color: Tuple[int, int, int]      # RGB
    secondary_color: Tuple[int, int, int]
    accent_color: Tuple[int, int, int]
    background_color: Tuple[int, int, int]
    text_color: Tuple[int, int, int]
    
    # 字體設定
    title_font: str
    body_font: str
    title_size: int
    body_size: int
    
    # 其他設定
    use_background_image: bool
    gradient_style: str  # 'none', 'horizontal', 'vertical'
    shadow_enabled: bool


# ==================== 風格預設 ====================
STYLE_PRESETS = {
    # 1. 現代科技風 (Modern Tech)
    'modern': StylePreset(
        name='Modern Tech',
        description='深藍色主題，簡潔現代感',
        primary_color=(30, 90, 160),        # 深藍
        secondary_color=(70, 130, 200),    # 淺藍
        accent_color=(0, 210, 200),        # 青綠
        background_color=(245, 245, 245),  # 淺灰
        text_color=(30, 30, 30),           # 深灰
        title_font='Arial',
        body_font='Arial',
        title_size=44,
        body_size=18,
        use_background_image=False,
        gradient_style='horizontal',
        shadow_enabled=True,
    ),
    
    # 2. 極簡風格 (Minimal Clean)
    'minimal': StylePreset(
        name='Minimal Clean',
        description='黑白灰色系，極簡設計',
        primary_color=(0, 0, 0),            # 黑
        secondary_color=(100, 100, 100),   # 灰
        accent_color=(255, 100, 0),        # 橙色點綴
        background_color=(255, 255, 255),  # 白
        text_color=(50, 50, 50),           # 深灰
        title_font='Arial',
        body_font='Arial',
        title_size=48,
        body_size=16,
        use_background_image=False,
        gradient_style='none',
        shadow_enabled=False,
    ),
    
    # 3. 企業正式風 (Corporate Professional)
    'corporate': StylePreset(
        name='Corporate Professional',
        description='深紅色主題，專業正式感',
        primary_color=(200, 40, 40),       # 深紅
        secondary_color=(240, 100, 100),  # 淺紅
        accent_color=(200, 150, 50),      # 金色
        background_color=(250, 250, 250), # 淺灰白
        text_color=(40, 40, 40),          # 深灰
        title_font='Calibri',
        body_font='Calibri',
        title_size=42,
        body_size=18,
        use_background_image=False,
        gradient_style='vertical',
        shadow_enabled=True,
    ),
    
    # 4. 創意藝術風 (Creative Artistic)
    'creative': StylePreset(
        name='Creative Artistic',
        description='紫色漸變，創意設計感',
        primary_color=(150, 80, 200),      # 紫色
        secondary_color=(100, 150, 255),  # 藍紫
        accent_color=(255, 200, 100),     # 溫暖黃
        background_color=(240, 235, 250), # 淺紫灰
        text_color=(60, 30, 80),          # 深紫灰
        title_font='Arial',
        body_font='Arial',
        title_size=44,
        body_size=18,
        use_background_image=False,
        gradient_style='horizontal',
        shadow_enabled=True,
    ),
    
    # 5. 清爽自然風 (Fresh Natural)
    'natural': StylePreset(
        name='Fresh Natural',
        description='綠色系主題，清爽自然',
        primary_color=(50, 140, 80),       # 深綠
        secondary_color=(100, 180, 120),  # 淺綠
        accent_color=(240, 150, 50),      # 溫暖橙
        background_color=(245, 250, 245), # 淺綠灰
        text_color=(30, 60, 30),          # 深綠灰
        title_font='Arial',
        body_font='Arial',
        title_size=42,
        body_size=18,
        use_background_image=False,
        gradient_style='vertical',
        shadow_enabled=False,
    ),
}


class PPTStyleConverter:
    """PPT 風格轉換器"""
    
    def __init__(self, input_file: str):
        """初始化轉換器
        
        Args:
            input_file: 輸入 PPT 檔案路徑
        """
        self.input_file = input_file
        self.output_dir = Path('./redesigned_ppts')
        self.output_dir.mkdir(exist_ok=True)
        
        try:
            self.prs = Presentation(input_file)
            print(f"✓ 成功加載 PPT: {input_file}")
            print(f"  - 投影片數: {len(self.prs.slides)}")
            print(f"  - 幻燈片尺寸: {self.prs.slide_width} x {self.prs.slide_height}")
        except Exception as e:
            print(f"✗ 無法加載 PPT: {e}")
            raise
    
    def apply_style_to_slide(self, slide, style: StylePreset):
        """將風格應用到單個投影片
        
        Args:
            slide: 投影片物件
            style: 風格預設
        """
        try:
            # 設定投影片背景
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*style.background_color)
            
            # 遍歷投影片中的所有形狀
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    
                    # 設定文字顏色
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(*style.text_color)
                        
                        # 標題處理
                        if hasattr(shape, 'name') and 'Title' in shape.name:
                            for run in paragraph.runs:
                                run.font.name = style.title_font
                                run.font.size = Pt(style.title_size)
                                run.font.bold = True
                                run.font.color.rgb = RGBColor(*style.primary_color)
                        # 正文處理
                        else:
                            for run in paragraph.runs:
                                run.font.name = style.body_font
                                run.font.size = Pt(style.body_size)
                
                # 設定形狀邊框和填充
                if shape.shape_type == 14:  # 文字框
                    if hasattr(shape, 'line'):
                        shape.line.color.rgb = RGBColor(*style.secondary_color)
                        shape.line.width = Pt(1)
                    
                    if hasattr(shape, 'fill'):
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = RGBColor(*style.accent_color)
                        shape.fill.transparency = 0.9  # 99% 透明
        
        except Exception as e:
            print(f"  ! 在處理形狀時出現警告: {e}")
    
    def redesign_with_style(self, style_name: str) -> str:
        """使用指定風格重新設計 PPT
        
        Args:
            style_name: 風格名稱 (必須在 STYLE_PRESETS 中)
            
        Returns:
            輸出檔案路徑
        """
        if style_name not in STYLE_PRESETS:
            raise ValueError(f"未知風格: {style_name}")
        
        style = STYLE_PRESETS[style_name]
        
        print(f"\n📝 應用風格: {style.name}")
        print(f"   描述: {style.description}")
        
        # 建立輸出演示文稿副本
        output_prs = Presentation(self.input_file)
        
        # 應用風格到所有投影片
        for idx, slide in enumerate(output_prs.slides):
            print(f"   處理投影片 {idx + 1}/{len(output_prs.slides)}...", end='\r')
            self.apply_style_to_slide(slide, style)
        
        # 生成輸出檔案名
        input_name = Path(self.input_file).stem
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_file = self.output_dir / f"{input_name}_{style_name}_{timestamp}.pptx"
        
        # 儲存檔案
        output_prs.save(str(output_file))
        print(f"\n✓ 完成: {style_name}")
        print(f"  儲存位置: {output_file}")
        
        return str(output_file)
    
    def batch_redesign(self, styles: List[str] = None) -> List[str]:
        """批量重新設計 PPT
        
        Args:
            styles: 風格列表 (None 表示使用所有風格)
            
        Returns:
            輸出檔案列表
        """
        if styles is None:
            styles = list(STYLE_PRESETS.keys())
        
        output_files = []
        print(f"\n🎨 開始批量重新設計...")
        print(f"   總計 {len(styles)} 種風格")
        print("-" * 60)
        
        for style_name in styles:
            try:
                output_file = self.redesign_with_style(style_name)
                output_files.append(output_file)
            except Exception as e:
                print(f"✗ 處理風格 {style_name} 失敗: {e}")
        
        print("-" * 60)
        print(f"✓ 完成所有轉換！共產生 {len(output_files)} 個檔案")
        
        return output_files
    
    def list_available_styles(self):
        """列出所有可用的風格"""
        print("\n📚 可用風格列表:")
        print("-" * 60)
        for style_name, style in STYLE_PRESETS.items():
            print(f"\n  {style_name.upper()}")
            print(f"    名稱: {style.name}")
            print(f"    描述: {style.description}")
            print(f"    主色: RGB{style.primary_color}")
            print(f"    字體: {style.title_font} / {style.body_font}")
        print("-" * 60)


def main():
    """命令行介面"""
    parser = argparse.ArgumentParser(
        description='PPT 風格自動重新設計工具',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog='''
範例用法:
  # 使用特定風格
  python ppt_style_converter.py input.pptx --styles modern minimal
  
  # 使用所有風格
  python ppt_style_converter.py input.pptx --all
  
  # 列出所有可用風格
  python ppt_style_converter.py --list
        '''
    )
    
    parser.add_argument('input', nargs='?', help='輸入 PPT 檔案路徑')
    parser.add_argument('--styles', nargs='+', help='指定風格 (空格分隔)')
    parser.add_argument('--all', action='store_true', help='使用所有風格')
    parser.add_argument('--list', action='store_true', help='列出所有可用風格')
    
    args = parser.parse_args()
    
    # 列出風格
    if args.list:
        converter = None
        temp_prs = Presentation()  # 臨時使用
        try:
            converter = PPTStyleConverter.__new__(PPTStyleConverter)
            converter.list_available_styles()
        except:
            print("使用方法: python ppt_style_converter.py input.pptx --styles modern minimal")
        return
    
    # 檢查輸入檔案
    if not args.input:
        parser.print_help()
        print("\n✗ 錯誤: 請指定輸入 PPT 檔案")
        sys.exit(1)
    
    if not os.path.exists(args.input):
        print(f"✗ 錯誤: 檔案不存在 - {args.input}")
        sys.exit(1)
    
    # 初始化轉換器
    converter = PPTStyleConverter(args.input)
    converter.list_available_styles()
    
    # 執行轉換
    if args.all:
        converter.batch_redesign()
    elif args.styles:
        converter.batch_redesign(args.styles)
    else:
        # 預設: 使用前 2 種風格
        converter.batch_redesign(['modern', 'minimal'])


if __name__ == '__main__':
    main()
